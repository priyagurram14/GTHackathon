import os
from pptx import Presentation
from pptx.util import Inches

def build_pptx(output_path, title, summary, narrative, charts=None, template_path=None):
    # Use template if available
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = 'Automated Insight Engine'

    # Narrative slide
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = 'Executive Summary'
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = narrative

    # Summary slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = 'Dataset Summary'
    if len(slide.placeholders) > 1:
        # summary may be dict; convert to readable lines
        if isinstance(summary, dict):
            lines = []
            lines.append(f"Rows: {summary.get('rows')}")
            lines.append(f"Cols: {summary.get('cols')}")
            for k,v in list(summary.get('missing',{}).items())[:8]:
                lines.append(f"Missing {k}: {v}")
            slide.placeholders[1].text = '\n'.join(lines)
        elif isinstance(summary, list):
            slide.placeholders[1].text = '\n'.join(summary)
        else:
            slide.placeholders[1].text = str(summary)

    # Charts
    if charts:
        for title_text, img_path in charts.items():
            # choose a layout that has a title and content
            layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            try:
                slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.4), width=Inches(8))
            except Exception as e:
                # if image can't be added, add the path as text
                body = slide.shapes.placeholders[1].text_frame
                body.text = f"Chart image missing: {img_path}\nError: {e}"

    # Save
    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
    prs.save(output_path)
